from pptx import Presentation

# Create a PowerPoint presentation object
prs = Presentation()

# Slide 1: Title Slide
slide_1 = prs.slides.add_slide(prs.slide_layouts[0])  # 0 is the layout for title slide
title = slide_1.shapes.title
subtitle = slide_1.placeholders[1]
title.text = "Atoms and Molecules"
subtitle.text = "Class: 9th\nSubject: Science\nPresented by: [Your Name]"

# Slide 2: Introduction
slide_2 = prs.slides.add_slide(prs.slide_layouts[1])  # 1 is the layout for a regular slide
title = slide_2.shapes.title
content = slide_2.shapes.placeholders[1]
title.text = "Introduction"
content.text = """
- Everything around us is made up of matter.
- Matter consists of tiny particles called atoms and molecules.
- Understanding atoms and molecules is crucial for studying chemistry.
- This chapter explores the fundamental concepts of atoms, molecules, and their interactions.
"""

# Slide 3: What is an Atom?
slide_3 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_3.shapes.title
content = slide_3.shapes.placeholders[1]
title.text = "What is an Atom?"
content.text = """
- The smallest unit of an element that retains its properties.
- Cannot be broken down further by chemical means.
- Composed of subatomic particles:
    - Protons (positively charged, found in the nucleus)
    - Neutrons (neutral charge, found in the nucleus)
    - Electrons (negatively charged, orbit around the nucleus)
- Examples: Hydrogen (H), Oxygen (O), Carbon (C)
- First proposed by John Dalton in his Atomic Theory.
"""

# Slide 4: Dalton’s Atomic Theory
slide_4 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_4.shapes.title
content = slide_4.shapes.placeholders[1]
title.text = "Dalton’s Atomic Theory"
content.text = """
- All matter is made of tiny, indivisible particles called atoms.
- Atoms of the same element are identical in mass and properties.
- Atoms of different elements differ in mass and properties.
- Atoms combine in whole-number ratios to form compounds.
- Chemical reactions involve the rearrangement of atoms without creation or destruction.
- Limitations:
    - Atoms are divisible (discovery of subatomic particles).
    - Isotopes exist (atoms of the same element with different masses).
"""

# Slide 5: Thomson’s Plum Pudding Model
slide_5 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_5.shapes.title
content = slide_5.shapes.placeholders[1]
title.text = "Thomson’s Plum Pudding Model"
content.text = """
- Proposed by J.J. Thomson in 1897.
- Atoms are made up of a positively charged sphere with negatively charged electrons embedded inside (like a plum pudding).
- Failed to explain the atomic nucleus and stability of electrons.
"""

# Slide 6: Rutherford’s Nuclear Model
slide_6 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_6.shapes.title
content = slide_6.shapes.placeholders[1]
title.text = "Rutherford’s Nuclear Model"
content.text = """
- Proposed by Ernest Rutherford in 1911 after the gold foil experiment.
- Observations:
    - Most alpha particles passed through gold foil (atom is mostly empty space).
    - Some deflected at small angles (suggested presence of a dense positive center).
    - A few bounced back (indicated a very small, dense, positively charged nucleus).
- Conclusions:
    - Atoms have a dense central nucleus containing protons.
    - Electrons revolve around the nucleus in empty space.
    - Could not explain the stability of electrons in orbits.
"""

# Slide 7: Bohr’s Atomic Model
slide_7 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_7.shapes.title
content = slide_7.shapes.placeholders[1]
title.text = "Bohr’s Atomic Model"
content.text = """
- Proposed by Niels Bohr in 1913.
- Postulates:
    - Electrons move in fixed orbits (energy levels) around the nucleus.
    - Electrons do not lose energy while revolving in these orbits.
    - Energy is absorbed or emitted when electrons jump between levels.
- Explained the stability of atoms and atomic spectra.
"""

# Slide 8: Modern Atomic Theory
slide_8 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_8.shapes.title
content = slide_8.shapes.placeholders[1]
title.text = "Modern Atomic Theory"
content.text = """
- Atoms are divisible and consist of protons, neutrons, and electrons.
- Electrons do not revolve in fixed orbits but in probabilistic regions called orbitals.
- The behavior of electrons follows quantum mechanics.
- Atoms combine based on valency (number of electrons lost, gained, or shared).
"""

# Continue creating slides for other topics similarly...

# Save the presentation
prs.save('Atoms_and_Molecules.pptx')

print("Presentation saved as Atoms_and_Molecules.pptx")
